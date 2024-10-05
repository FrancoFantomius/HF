import fitz  # PyMuPDF
import docx
import pandas as pd
from pandas.plotting import table
import matplotlib
matplotlib.use('Agg')  # Use a non-interactive backend
import matplotlib.pyplot as plt
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io



## Function to generate an image from a pdf path
def pdf_page_to_image(pdf_path, page_number, dpi=200):
    """
    Convert a specific PDF page to an image using PyMuPDF and Pillow.

    :param pdf_path: Path to the PDF file.
    :param page_number: The page number to convert (0-indexed).
    :param dpi: Resolution of the output image.
    :return: PIL.Image object of the rendered page.
    """
    # Open the PDF file
    pdf_document = fitz.open(pdf_path)
    
    # Ensure the page_number is valid
    if page_number < 0 or page_number >= pdf_document.page_count:
        raise ValueError("Invalid page number")
    
    # Select the page
    page = pdf_document.load_page(page_number)  # 0-indexed
    
    # Convert PDF page to a pixmap (image) at the specified DPI
    # DPI scaling is controlled by the zoom matrix
    zoom = dpi / 72  # 72 is the default DPI of PDF
    mat = fitz.Matrix(zoom, zoom)
    
    pix = page.get_pixmap(matrix=mat)
    
    # Convert the pixmap to a byte buffer
    img_bytes = pix.tobytes("png")
    
    # Use Pillow to open the image from bytes
    img = Image.open(io.BytesIO(img_bytes))
    
    return img


#Function to render all the word document in a single image
def docx_to_image(docx_path):
    doc = docx.Document(docx_path)
    
    # Font for rendering text
    font = ImageFont.load_default()
    
    text = ""

    # Convert each paragraph to an image
    for para in doc.paragraphs:
        text += para.text + "\n"

    img = Image.new('RGB', (800, 600), color='white')
    d = ImageDraw.Draw(img)
    d.text((10, 10), text, font=font, fill='black')
    
    return img


def txt_to_image(txt_path):
    with open(txt_path, 'r') as file:
        text = file.read()
    
    font = ImageFont.load_default()
    
    img = Image.new('RGB', (800, 100), color='white')
    d = ImageDraw.Draw(img)
    d.text((10, 10), text, font=font, fill='black')
    
    return img


def excel_sheet_to_images(excel_path, sheet_nmb):
    # Step 1: Read the specified sheet from the Excel file
    excel_file = pd.ExcelFile(excel_path)
    sheet_name = excel_file.sheet_names[int(sheet_nmb)]

    df = pd.read_excel(excel_path, sheet_name)
    # Step 2: Create a Matplotlib figure to render the DataFrame
    fig, ax = plt.subplots(figsize=(10, 6))  # Adjust the figure size as needed
    ax.axis('tight')
    ax.axis('off')
    
    # Step 3: Create a table from the DataFrame and render it in the figure
    tbl = table(ax, df, loc='center', cellLoc='center', colWidths=[0.1] * len(df.columns))
    
    # Step 4: Create an in-memory bytes buffer for the image
    img_buffer = io.BytesIO()
    
    # Save the figure to the in-memory buffer
    plt.savefig(img_buffer, format='png', dpi=300)
    
    # Close the plot to release memory
    plt.close(fig)
    
    # Seek to the start of the buffer so Flask can read it
    img_buffer.seek(0)
    
    # Return the in-memory image buffer
    return img_buffer

def csv_to_images(csv_path):
    df = pd.read_csv(csv_path)
    
    # Convert the CSV data to an image using pandas
    fig, ax = plt.subplots(figsize=(8, 6))
    ax.axis('tight')
    ax.axis('off')
    ax.table(cellText=df.values, colLabels=df.columns, loc='center')
    
    buf = io.BytesIO()
    plt.savefig(buf, format='png')
    buf.seek(0)
    img = Image.open(buf)
    buf.close()
    
    return img


def pptx_to_images(pptx_path, slide_nmb):
    prs = Presentation(pptx_path)
    
    slide = prs.slides[int(slide_nmb)]
    img = Image.new('RGB', (1280, 720), color='white')
    d = ImageDraw.Draw(img)
    
    # Extract the slide text
    text = '\n'.join([shape.text for shape in slide.shapes if hasattr(shape, 'text')])
    
    # Draw the slide text onto the image
    d.text((10, 10), text, fill='black')

    return img
