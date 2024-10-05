import os
import docx
import pptx
import openpyxl
import fitz  # PyMuPDF
import chardet

def search_text_in_file(text, file_path):
    """
    Searches for a specific text within a list of files.
    Supported file formats: .docx, .pptx, .xlsx, .pdf, .rtf, .txt
    
    :param text: The text to search for.
    :param file_paths: List of file paths.
    :return: Dictionary with file paths and boolean indicating if text was found.
    """
    extension = os.path.splitext(file_path)[1].lower()
    
    if extension == '.docx':
        found = search_in_docx(file_path, text)
    elif extension == '.pptx':
        found = search_in_pptx(file_path, text)
    elif extension == '.xlsx':
        found = search_in_xlsx(file_path, text)
    elif extension == '.pdf':
        found = search_in_pdf(file_path, text)
    elif extension == '.txt':
        found = search_in_txt(file_path, text)
    else:
        found = False
    
    return found


def search_in_docx(file_path, search_text):
    doc = docx.Document(file_path)
    for paragraph in doc.paragraphs:
        if search_text.lower() in paragraph.text.lower():
            return True
    return False


def search_in_pptx(file_path, search_text):
    ppt = pptx.Presentation(file_path)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and search_text.lower() in shape.text.lower():
                return True
    return False


def search_in_xlsx(file_path, search_text):
    wb = openpyxl.load_workbook(file_path)
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and search_text.lower() in str(cell.value).lower():
                    return True
    return False


def search_in_pdf(file_path, search_text):
    document = fitz.open(file_path)
    for page in document:
        text = page.get_text()
        if search_text.lower() in text.lower():
            return True
    return False


def search_in_txt(file_path, search_text):
    with open(file_path, 'rb') as f:
        raw_data = f.read()
        encoding = chardet.detect(raw_data)['encoding']
    
    with open(file_path, 'r', encoding=encoding) as f:
        content = f.read()
        if search_text.lower() in content.lower():
            return True
    return False
'''
# Example Usage:
file_list = ['file1.docx', 'file2.pptx', 'file3.xlsx', 'file4.pdf', 'file5.rtf', 'file6.txt']
search_term = 'your search text'

result = search_text_in_files(search_term, file_list)
for file, found in result.items():
    print(f"File: {file}, Found: {found}")
'''